"""
ORIGIN: Manual creation to convert PRESENTATION.md into PowerPoint presentation.pptx
PURPOSE: Parse the master slide deck markdown and generate a 16:9 dark-mode Teal/Sky
         PowerPoint (.pptx). Layout engine supports a subtitle band, two-column regions
         (via a '+++' marker), full-width or in-column ASCII diagrams (```-fenced) that
         auto-fit their box, larger body fonts that fill the page, and speaker notes
         (everything after a 'NOTES:' marker) injected into the slide's notes pane.
DESTINATION: Output saved as presentation.pptx in repository root (git-ignored, local only).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (Teal & Sky tech theme, dark) ----
BG_COLOR = RGBColor(9, 13, 22)          # Deep slate
CARD_BG = RGBColor(15, 23, 42)          # Slate 900
TITLE_COLOR = RGBColor(56, 189, 248)    # Sky 400
SUBTITLE_COLOR = RGBColor(45, 212, 191) # Teal 400
TEXT_COLOR = RGBColor(241, 245, 249)    # Slate 100
ACCENT_COLOR = RGBColor(13, 148, 136)   # Teal 600
MUTED_COLOR = RGBColor(148, 163, 184)   # Slate 400
CODE_COLOR = RGBColor(226, 232, 240)

CONTENT_BOTTOM = 7.18   # inches — keep everything above this
LEFT_X, COL_W = 0.55, 5.98
RIGHT_X = 6.82


def fit_font(longest_chars, n_lines, w_in, h_in, cap=13.0, floor=7.0):
    """Largest Consolas size that fits `longest_chars` wide and `n_lines` tall in the box."""
    by_w = 72.0 * w_in / (0.62 * max(longest_chars, 1))   # Consolas ~0.55em + safety margin
    by_h = 72.0 * h_in / (1.18 * max(n_lines, 1))
    return max(floor, min(cap, by_w, by_h))


def _zero_margins(tf):
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0


def add_diagram(slide, code_lines, x, y, w, h):
    """Rounded card with a monospace ASCII diagram, auto-sized to fit."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_COLOR
    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.14),
                                  Inches(w - 0.4), Inches(h - 0.28))
    tf = tb.text_frame
    tf.word_wrap = False
    _zero_margins(tf)
    longest = max((len(l) for l in code_lines), default=1)
    fs = fit_font(longest, len(code_lines), w - 0.45, h - 0.3, cap=13.0, floor=6.5)
    for i, cl in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = cl
        p.font.name = "Consolas"
        p.font.size = Pt(fs)
        p.font.color.rgb = CODE_COLOR
        p.line_spacing = 1.0


def _add_formatted_text(tf, raw_text, font_size=14.0):
    lines = raw_text.split("\n")
    first = True
    for line in lines:
        sline = line.strip()
        if not sline:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if sline.startswith("#### "):
            p.text = sline[5:].strip().replace("**", "")
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size + 2)
            p.font.bold = True
            p.font.color.rgb = SUBTITLE_COLOR
            p.space_after = Pt(4)
            p.space_before = Pt(4)
        elif sline.startswith("### "):
            p.text = sline[4:].strip().replace("**", "")
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size + 4)
            p.font.bold = True
            p.font.color.rgb = SUBTITLE_COLOR
            p.space_after = Pt(5)
        elif sline.startswith("- ") or sline.startswith("* "):
            p.text = "▸  " + sline[2:].replace("**", "").replace("[", "").replace("]", "")
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(6)
            p.line_spacing = 1.12
        elif sline.startswith("> "):
            p.text = sline[2:].strip().replace("**", "")
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.italic = True
            p.font.color.rgb = MUTED_COLOR
            p.space_after = Pt(5)
            p.line_spacing = 1.15
        else:
            p.text = sline.replace("**", "").replace("[", "").replace("]", "")
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(6)
            p.line_spacing = 1.15


def add_text_block(slide, text, x, y, w, h, base=14.0):
    """Bullets/paragraphs, font shrinking a touch when the block is dense."""
    n = len([l for l in text.split("\n") if l.strip()])
    fs = base if n <= 7 else (base - 1.5 if n <= 10 else base - 3)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    _add_formatted_text(tf, text, font_size=fs)


def render_region(slide, content, x, y, w, h):
    """Render one region: bullets, a full diagram, or short text above a diagram."""
    content = content.strip()
    if "```" in content:
        parts = content.split("```")
        pre = parts[0].strip()
        code = parts[1].strip() if len(parts) > 1 else ""
        code_lines = code.split("\n")
        if code_lines and code_lines[0].strip().isalpha():
            code_lines = code_lines[1:]
        cy, ch = y, h
        if len(pre) > 15:
            pre_n = len([l for l in pre.split("\n") if l.strip()])
            th = min(h * 0.44, 0.34 * pre_n + 0.25)
            add_text_block(slide, pre, x, cy, w, th, base=14)
            cy += th + 0.12
            ch = h - (cy - y)
        add_diagram(slide, code_lines, x, cy, w, ch)
    else:
        add_text_block(slide, content, x, y, w, h, base=14.5)


def create_presentation():
    md_path = "PRESENTATION.md"
    if not os.path.exists(md_path):
        print(f"Error: {md_path} not found.")
        return
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    raw_slides = [s.strip() for s in content.split("\n---\n") if s.strip()]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, slide_md in enumerate(raw_slides):
        slide = prs.slides.add_slide(blank)

        # background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # title + body/notes split
        title_text = f"Slide {idx + 1}"
        body_lines = []
        for line in slide_md.split("\n"):
            if line.startswith("## Slide"):
                title_text = line.lstrip("#").strip()
            elif line.startswith("# ") and idx == 0:
                title_text = line.lstrip("#").strip()
            else:
                body_lines.append(line)
        body_text = "\n".join(body_lines).strip()

        notes_text = ""
        if "\nNOTES:" in body_text:
            body_text, notes_text = body_text.split("\nNOTES:", 1)
            body_text, notes_text = body_text.strip(), notes_text.strip()
        elif body_text.startswith("NOTES:"):
            notes_text, body_text = body_text[6:].strip(), ""
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

        # header
        hb = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.13), Inches(0.7))
        tf = hb.text_frame
        tf.word_wrap = True
        _zero_margins(tf)
        ph = tf.paragraphs[0]
        ph.text = title_text
        ph.font.name = "Segoe UI"
        ph.font.size = Pt(23)
        ph.font.bold = True
        ph.font.color.rgb = TITLE_COLOR
        # accent underline
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.08), Inches(12.13), Inches(0.03))
        ln.fill.solid()
        ln.fill.fore_color.rgb = ACCENT_COLOR
        ln.line.fill.background()

        # optional subtitle band (first '### ' line)
        content_top = 1.32
        rest = body_text
        blines = body_text.split("\n")
        if blines and blines[0].strip().startswith("### "):
            sub = blines[0].strip()[4:].replace("**", "").strip()
            rest = "\n".join(blines[1:]).strip()
            sb = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.13), Inches(0.66))
            tfs = sb.text_frame
            tfs.word_wrap = True
            _zero_margins(tfs)
            ps = tfs.paragraphs[0]
            ps.text = sub
            ps.font.name = "Segoe UI"
            ps.font.size = Pt(18)
            ps.font.bold = True
            ps.font.color.rgb = SUBTITLE_COLOR
            content_top = 2.02

        region_h = CONTENT_BOTTOM - content_top

        # columns via '+++' else full width
        if "\n+++\n" in rest:
            left, right = rest.split("\n+++\n", 1)
            render_region(slide, left, LEFT_X, content_top, COL_W, region_h)
            render_region(slide, right, RIGHT_X, content_top, COL_W, region_h)
        else:
            render_region(slide, rest, LEFT_X, content_top, 12.23, region_h)

    out_path = os.environ.get("PPTX_OUT", "presentation.pptx")
    prs.save(out_path)
    print(f"Successfully generated {out_path} with {len(raw_slides)} slides (Teal & Sky Theme).")


if __name__ == "__main__":
    create_presentation()
